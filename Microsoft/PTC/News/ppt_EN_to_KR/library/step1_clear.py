"""
step1_clear.py — STEP 1: 원본 PPTX 복사 + 슬라이드 Shape 클리어

처리 흐름:
  1. eng/{파일명}.pptx → kr/{파일명}_KO.pptx 복사
  2. kr/ 작업본의 모든 슬라이드에서 Shape 제거 (레이아웃·마스터·테마 유지)
  3. 슬라이드 노트도 클리어
  4. 작업본 저장
"""
import os
import shutil

from pptx import Presentation

from library.config import Config


def run(pptx_path: str, cfg: Config, stem: str) -> None:
    """
    원본을 kr/ 작업본으로 복사한다.
    (텍스트 in-place 치환 방식으로 변경되어 Shape 클리어는 비활성화)

    Parameters
    ----------
    pptx_path : 원본 .pptx 절대 경로 (eng/ 또는 temp/)
    cfg       : Config 인스턴스 (kr_dir 경로 사용)
    stem      : 확장자 없는 파일명
    """
    kr_path = _kr_path(cfg, stem)
    _copy(pptx_path, kr_path)
    # _clear_shapes(kr_path)  # 비활성화: 원본 Shape 유지하고 텍스트만 in-place 치환
    print(f"[STEP 1] 완료 → {os.path.basename(kr_path)}")


# ──────────────────────────────────────────────
# 내부 구현
# ──────────────────────────────────────────────

def _kr_path(cfg: Config, stem: str) -> str:
    os.makedirs(cfg.kr_dir, exist_ok=True)
    return os.path.join(cfg.kr_dir, f"{stem}_KO.pptx")


def _copy(src: str, dst: str) -> None:
    shutil.copy2(src, dst)
    print(f"  복사: {os.path.basename(src)} → {os.path.basename(dst)}")


def _clear_shapes(kr_path: str) -> None:
    prs = Presentation(kr_path)
    for idx, slide in enumerate(prs.slides, start=1):
        removed = _remove_shapes(slide)
        _clear_notes(slide)
        if removed:
            print(f"  slide_{idx}: {removed}개 Shape 제거")
    prs.save(kr_path)


def _remove_shapes(slide) -> int:
    sp_tree = slide.shapes._spTree
    _SHAPE_TAGS = {"sp", "pic", "graphicFrame", "grpSp", "cxnSp"}
    to_remove = [
        child for child in list(sp_tree)
        if _tag_name(child) in _SHAPE_TAGS
    ]
    for elem in to_remove:
        sp_tree.remove(elem)
    return len(to_remove)


def _clear_notes(slide) -> None:
    try:
        slide.notes_slide.notes_text_frame.text = ""
    except Exception:
        pass


def _tag_name(elem) -> str:
    tag = elem.tag
    return tag.split("}")[-1] if "}" in tag else tag
