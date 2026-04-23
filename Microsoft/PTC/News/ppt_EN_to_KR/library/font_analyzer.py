"""
font_analyzer.py — STEP 2 보조: 슬라이드별 영어 폰트 목록 수집 + 한글 폰트 대응
신규 폰트 발견 시 규칙 기반으로 매핑하고, 실패 시 기본값(맑은 고딕)으로 폴백한다.
LLM 호출은 하지 않는다 (font.json 규칙 기반으로 충분).
"""
import json
import os

from openai import AzureOpenAI
from pptx import Presentation


# 규칙 기반 폰트 매핑: 폰트명(소문자)에 키워드가 포함되면 해당 한글 폰트 적용
# Pretendard 계열을 기본으로 사용 (모던 디자인, 무료 한글 폰트)
_RULE_BASED: list = [
    (["light", "thin"],                                       "Pretendard Light"),
    (["consolas", "courier", "mono", "code"],                 "D2Coding"),
    (["times", "georgia", "serif", "batang", "garamond"],     "본명조"),
    (["black", "heavy", "extrabold"],                          "Pretendard ExtraBold"),
    (["bold", "semibold", "demi"],                             "Pretendard SemiBold"),
    (["arial", "helvetica", "gothic"],                         "Pretendard"),
    (["segoe", "calibri", "tahoma", "verdana", "trebuchet"],   "Pretendard"),
]
_DEFAULT_KR_FONT = "Pretendard"


def analyze(pptx_path: str, work_dir: str, filename_stem: str,
            font_map_path: str, llm_client: AzureOpenAI, model: str) -> None:
    """
    원본 PPTX에서 슬라이드별 폰트를 수집하고 slide_{N}_font.json을 저장한다.
    신규 폰트는 규칙 기반으로 매핑하여 font.json에 추가한다. (LLM 미사용)
    """
    prs = Presentation(pptx_path)
    font_map = _load_font_map(font_map_path)
    out_dir = os.path.join(work_dir, "components_en", filename_stem)
    os.makedirs(out_dir, exist_ok=True)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_fonts = _collect_slide_fonts(slide)
        new_fonts = [f for f in slide_fonts if f and f not in font_map]
        if new_fonts:
            updates = _map_fonts_by_rule(new_fonts)
            font_map.update(updates)
            _save_font_map(font_map_path, font_map)

        # slide_{N}_font_en.json: 이 슬라이드에서 사용된 폰트만
        slide_font_map = {
            f: font_map.get(f, _DEFAULT_KR_FONT)
            for f in slide_fonts if f
        }
        slide_font_map["__default__"] = font_map.get("__default__", _DEFAULT_KR_FONT)

        out_path = os.path.join(out_dir, f"slide_{slide_idx}_font_en.json")
        with open(out_path, "w", encoding="utf-8") as fp:
            json.dump(slide_font_map, fp, ensure_ascii=False, indent=2)

    print(f"[font_analyzer] {filename_stem}: 폰트 분석 완료")


def _collect_slide_fonts(slide) -> set:
    fonts = set()
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    fonts.add(run.font.name)
        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    try:
                        for para in cell.text_frame.paragraphs:
                            for run in para.runs:
                                fonts.add(run.font.name)
                    except Exception:
                        pass
    return fonts


def _map_fonts_by_rule(fonts: list) -> dict:
    """폰트명 키워드 매칭으로 한글 폰트를 결정한다."""
    result = {}
    for font in fonts:
        lower = font.lower()
        kr = _DEFAULT_KR_FONT
        for keywords, mapped in _RULE_BASED:
            if any(kw in lower for kw in keywords):
                kr = mapped
                break
        result[font] = kr
        print(f"  [font_analyzer] 신규 폰트 매핑: {font!r} → {kr}")
    return result


def _load_font_map(path: str) -> dict:
    if os.path.exists(path):
        with open(path, encoding="utf-8") as f:
            return json.load(f)
    return {"__default__": "맑은 고딕"}


def _save_font_map(path: str, data: dict) -> None:
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
