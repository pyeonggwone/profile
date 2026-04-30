"""
doc_generator.py — STEP 4: 설명자료 PPTX 생성 (기본 실행 제외)
template_guide.pptx 없으면 자동 스킵.
"""
import json
import os

from openai import AzureOpenAI
from pptx import Presentation
from pptx.util import Pt


def generate(
    template_path: str,
    comp_dir: str,
    filename_stem: str,
    kr_dir: str,
    llm_client: AzureOpenAI,
    model: str,
    total_slides: int,
) -> None:
    """
    template_guide.pptx 기반으로 설명자료 PPTX를 생성한다.
    template_guide.pptx 없으면 스킵.
    """
    if not os.path.exists(template_path):
        print("[doc_generator] template_guide.pptx 없음 → STEP 4 스킵")
        return

    # 1) 템플릿 레이아웃 정보 수집
    prs_template = Presentation(template_path)
    layouts_info = [
        {"index": i, "name": layout.name}
        for i, layout in enumerate(prs_template.slide_layouts)
    ]

    # 2) STEP 2 component JSON 수집 (슬라이드별 제목/내용 요약)
    slides_summary = []
    for slide_idx in range(1, total_slides + 1):
        comp_path = os.path.join(comp_dir, filename_stem, f"slide_{slide_idx}_component.json")
        if not os.path.exists(comp_path):
            continue
        with open(comp_path, encoding="utf-8") as f:
            comp = json.load(f)
        texts = []
        for tb in comp.get("text_boxes", []):
            for para in tb.get("paragraphs", []):
                t = para.get("text", "").strip()
                if t:
                    texts.append(t)
        slides_summary.append({"slide_num": slide_idx, "texts": texts[:5]})  # 최대 5개 텍스트

    # 3) LLM: 설명자료 슬라이드 구성 JSON 생성
    prompt = (
        "아래 원본 PPTX 슬라이드 내용을 바탕으로 설명자료 슬라이드를 구성하라.\n"
        f"사용 가능한 레이아웃: {json.dumps(layouts_info, ensure_ascii=False)}\n\n"
        f"슬라이드 내용 요약:\n{json.dumps(slides_summary, ensure_ascii=False)}\n\n"
        "결과 형식 (JSON):\n"
        '{"slides": [{"layout_index": 1, "title": "개요", "body": "내용..."}, ...]}'
    )
    response = llm_client.chat.completions.create(
        model=model,
        messages=[{"role": "user", "content": prompt}],
        response_format={"type": "json_object"},
    )
    guide_data: dict = json.loads(response.choices[0].message.content)

    # 4) python-pptx로 슬라이드 생성
    prs_out = Presentation(template_path)
    # 기존 슬라이드 제거
    for _ in range(len(prs_out.slides)):
        slide_id = prs_out.slides._sldIdLst[0]
        prs_out.slides._sldIdLst.remove(slide_id)

    for slide_info in guide_data.get("slides", []):
        layout_idx = slide_info.get("layout_index", 0)
        layout = prs_out.slide_layouts[min(layout_idx, len(prs_out.slide_layouts) - 1)]
        slide = prs_out.slides.add_slide(layout)

        title_placeholder = None
        body_placeholder = None
        for ph in slide.placeholders:
            if ph.placeholder_format.idx == 0:
                title_placeholder = ph
            elif ph.placeholder_format.idx == 1:
                body_placeholder = ph

        if title_placeholder:
            title_placeholder.text = slide_info.get("title", "")
        if body_placeholder:
            body_placeholder.text = slide_info.get("body", "")

    out_path = os.path.join(kr_dir, f"{filename_stem}_GUIDE.pptx")
    prs_out.save(out_path)
    print(f"[doc_generator] 설명자료 생성 완료 → {out_path}")
