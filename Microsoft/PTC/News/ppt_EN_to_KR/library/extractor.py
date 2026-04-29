"""
extractor.py — STEP 2: 슬라이드별 컴포넌트 직렬화
원본 PPTX에서 슬라이드별 Shape를 순회하여 상태를 추출하고
work/components/{파일명}/slide_{N}_component.json 에 저장한다.
"""
import json
import os

from typing import Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt


# Shape type 상수 (python-pptx enum 값)
_TABLE_TYPE = 19       # MSO_SHAPE_TYPE.TABLE
_PICTURE_TYPE = 13     # MSO_SHAPE_TYPE.LINKED_PICTURE / PICTURE
_MEDIA_TYPE = 16       # MSO_SHAPE_TYPE.MEDIA


def extract(pptx_path: str, work_dir: str, filename_stem: str) -> None:
    """
    원본 PPTX에서 슬라이드별 컴포넌트를 추출하여 JSON으로 저장한다.
    이미지는 work/img/{filename_stem}/slide_{N}/*.jpg 에 저장된다.
    """
    prs = Presentation(pptx_path)
    out_dir = os.path.join(work_dir, "components_en", filename_stem)
    img_dir = os.path.join(work_dir, "img", filename_stem)
    os.makedirs(out_dir, exist_ok=True)

    for slide_idx, slide in enumerate(prs.slides, start=1):
        data = _extract_slide(slide, slide_idx, img_dir, filename_stem)
        out_path = os.path.join(out_dir, f"slide_{slide_idx}_component_en.json")
        with open(out_path, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)

    print(f"[extractor] {filename_stem}: {len(prs.slides)} 슬라이드 추출 완료")


def _extract_slide(slide, slide_num: int, img_dir: str, filename_stem: str) -> dict:
    data: dict = {"slide_num": slide_num}
    text_boxes = []
    images = []
    tables = []
    smartarts = []
    charts = []

    slide_img_dir = os.path.join(img_dir, f"slide_{slide_num}")

    for shape in slide.shapes:
        shape_id = f"s{slide_num}_shape{shape.shape_id}"

        # 표(Table)
        if shape.shape_type == _TABLE_TYPE or shape.has_table:
            tables.append(_extract_table(shape, shape_id))
            continue

        # 이미지(Picture)
        if shape.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.LINKED_PICTURE):
            entry = _extract_image(shape, shape_id, slide_img_dir)
            if entry:
                images.append(entry)
            continue

        # 차트
        if shape.has_chart:
            charts.append({
                "id": shape_id,
                "left": shape.left, "top": shape.top,
                "width": shape.width, "height": shape.height,
                "shape_type": int(shape.shape_type),
            })
            continue

        # 텍스트 박스 / 제목 등
        if shape.has_text_frame:
            text_boxes.append(_extract_text(shape, shape_id))
            continue

        # SmartArt (shape_type == 6 또는 GROUP)
        if shape.shape_type in (MSO_SHAPE_TYPE.GROUP, 6):
            smartarts.append({
                "id": shape_id,
                "left": shape.left, "top": shape.top,
                "width": shape.width, "height": shape.height,
                "shape_type": int(shape.shape_type),
            })
            continue

        # 그 외 처리 불가 Shape
        smartarts.append({
            "id": shape_id,
            "left": shape.left, "top": shape.top,
            "width": shape.width, "height": shape.height,
            "shape_type": int(shape.shape_type),
        })

    # 비어있는 키는 제거
    if text_boxes:
        data["text_boxes"] = text_boxes
    if images:
        data["images"] = images
    if tables:
        data["tables"] = tables
    if smartarts:
        data["smartarts"] = smartarts
    if charts:
        data["charts"] = charts

    # 슬라이드 노트
    notes_text = _extract_notes(slide)
    if notes_text:
        data["notes"] = notes_text

    return data


def _extract_text(shape, shape_id: str) -> dict:
    paragraphs = []
    for para in shape.text_frame.paragraphs:
        runs = []
        for run in para.runs:
            font = run.font
            runs.append({
                "text": run.text,
                "font": font.name,
                "bold": font.bold,
                "italic": font.italic,
                "underline": font.underline,
                "size": int(font.size / 12700) if font.size else None,  # EMU → pt
                "color": _color_hex(font.color),
            })
        # 런이 없는 단락도 텍스트 보존
        para_text = para.text
        paragraphs.append({"runs": runs, "text": para_text})

    return {
        "id": shape_id,
        "left": shape.left, "top": shape.top,
        "width": shape.width, "height": shape.height,
        "paragraphs": paragraphs,
    }


def _extract_image(shape, shape_id: str, slide_img_dir: str) -> Optional[dict]:
    try:
        image = shape.image
        os.makedirs(slide_img_dir, exist_ok=True)
        ext = image.content_type.split("/")[-1]
        if ext == "jpeg":
            ext = "jpg"
        img_count = len(os.listdir(slide_img_dir)) + 1
        img_filename = f"image_{img_count}.{ext}"
        img_path = os.path.join(slide_img_dir, img_filename)
        with open(img_path, "wb") as f:
            f.write(image.blob)
        return {
            "id": shape_id,
            "left": shape.left, "top": shape.top,
            "width": shape.width, "height": shape.height,
            "img_path": img_path.replace("\\", "/"),
        }
    except Exception as e:
        print(f"[extractor] 이미지 추출 실패 ({shape_id}): {e}")
        return None


def _extract_table(shape, shape_id: str) -> dict:
    table = shape.table
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            tf = cell.text_frame
            cell_text = tf.text if tf else ""
            font_info = {}
            # 첫 번째 Run 서식 참조
            try:
                run = tf.paragraphs[0].runs[0]
                font = run.font
                font_info = {
                    "font": font.name,
                    "bold": font.bold,
                    "size": int(font.size / 12700) if font.size else None,
                    "color": _color_hex(font.color),
                }
            except (IndexError, AttributeError):
                pass
            cells.append({"text": cell_text, **font_info})
        rows.append(cells)
    return {
        "id": shape_id,
        "left": shape.left, "top": shape.top,
        "width": shape.width, "height": shape.height,
        "rows": rows,
    }


def _extract_notes(slide) -> str:
    try:
        notes_slide = slide.notes_slide
        return notes_slide.notes_text_frame.text.strip()
    except Exception:
        return ""


def _color_hex(color_obj) -> Optional[str]:
    try:
        rgb = color_obj.rgb
        return f"#{rgb}"
    except Exception:
        return None
